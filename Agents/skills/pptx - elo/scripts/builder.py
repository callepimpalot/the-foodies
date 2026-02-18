from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import json

class SlideBuilder:
    def __init__(self, brand_config):
        """Initializes the presentation with a standard 16:9 aspect ratio."""
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(5.625)
        self.brand = brand_config

    def apply_text_style(self, shape, text_data, is_header=False):
        """Applies typography and color rules from the separate Style Contract."""
        tf = shape.text_frame
        tf.word_wrap = True
        
        # Ensure at least one paragraph exists
        if not tf.paragraphs:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
            
        p.text = text_data
        
        # Apply style from config to the first run
        if p.runs:
            run = p.runs[0]
            run.font.name = self.brand.get('header_font' if is_header else 'body_font', 'Arial')
            run.font.size = Pt(self.brand.get('header_size' if is_header else 'body_size', 18))
            color_hex = self.brand.get('primary_color', '000000').lstrip('#')
            run.font.color.rgb = RGBColor.from_string(color_hex)

    def create_asymmetric_slide(self, title, bullets):
        """Creates an 'Anti-Boring' slide with an asymmetric 40/60 split."""
        # Use a blank layout (usually index 6)
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6]) 
        
        # Left Rail (40% width for Title)
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(3.5), Inches(4.5))
        self.apply_text_style(title_shape, title, is_header=True)
        
        # Right Rail (60% width for Body Content)
        body_shape = slide.shapes.add_textbox(Inches(4.5), Inches(0.5), Inches(5), Inches(4.5))
        
        # Join bullets with newlines and check for bottom margin (0.5")
        body_text = "\n".join(bullets)
        self.apply_text_style(body_shape, body_text)

    def save(self, filename):
        self.prs.save(filename)

def build_from_spec(spec_path, brand_path, output_name):
    """Bridge function to execute the build from your Execution Packet files."""
    with open(spec_path, 'r') as f:
        spec = json.load(f)
    with open(brand_path, 'r') as f:
        # Assuming brand-style.md is formatted as a simple JSON or dict
        brand = json.load(f) 
        
    builder = SlideBuilder(brand)
    for slide_id, content in spec.items():
        title = content.get('title', 'No Title')
        # Extract text from the SCS 'paragraphs' structure
        bullets = [p.get('text', '') for p in content.get('paragraphs', [])]
        builder.create_asymmetric_slide(title, bullets)
    
    builder.save(output_name)